#!/usr/bin/env python3
"""
scripts/ingest_brand_assets.py — Descarga brand assets de OneDrive/SharePoint
y los ingesta al RAG del Agent Designer (namespace=brand, tenant=reamerica).

Variables de entorno requeridas (vault o .env):
  ONEDRIVE_REAMERICA_USER    — email del usuario cuyo drive contiene las carpetas
  ONEDRIVE_REAMERICA_SITE_ID — (opcional) ID del site SharePoint si aplica
  OUTLOOK_REAMERICA_CLIENT_ID / CLIENT_SECRET / TENANT_ID — ya en vault

Uso:
  python scripts/ingest_brand_assets.py --user mromanelli@reamerica-re.com
  python scripts/ingest_brand_assets.py --user mromanelli@reamerica-re.com --site <SITE_ID>
  python scripts/ingest_brand_assets.py --list-sites  # ver sites SharePoint disponibles
"""
import argparse
import logging
import os
import sys
from datetime import datetime
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(name)s %(message)s",
)
log = logging.getLogger(__name__)

DOWNLOAD_DIR = Path(__file__).parent.parent / "downloads" / "onedrive"
BRAND_FOLDERS = ["marketing", "identidad de marca"]
BRAND_NAMESPACE = "brand"
BRAND_TENANT = "reamerica"

TEXT_EXTENSIONS = {".pdf", ".txt", ".md", ".pptx"}
VISUAL_EXTENSIONS = {".png", ".jpg", ".jpeg", ".svg", ".ai", ".eps",
                     ".ase", ".gif", ".webp", ".tif", ".tiff"}


# ── Extracción de texto ────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = [p.extract_text() or "" for p in reader.pages]
    return "\n\n".join(t.strip() for t in pages if t.strip())


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        return "\n".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, "text")
        )
    except ImportError:
        log.warning("python-pptx no instalado, skip pptx")
        return ""


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        if ext in (".txt", ".md"):
            return path.read_text(errors="ignore")
        if ext == ".pptx":
            return _extract_pptx(path)
    except Exception as e:
        log.warning(f"Error extrayendo texto de {path.name}: {e}")
    return ""


# ── Ingesta RAG ────────────────────────────────────────────────────────────────

def ingest_file(path: Path, folder_origin: str) -> dict:
    """Ingesta un archivo al RAG. Devuelve dict con resultado."""
    from modules.rag_kb import ingest

    ext = path.suffix.lower()
    result = {"file": str(path), "name": path.name, "folder": folder_origin,
              "ext": ext, "status": "skip", "chunks": 0}

    metadata = {
        "source_folder": folder_origin,
        "filename": path.name,
        "extension": ext,
        "ingested_at": datetime.utcnow().isoformat(),
        "type": "brand_asset",
    }

    if ext in TEXT_EXTENSIONS:
        text = extract_text(path)
        if not text.strip():
            result["status"] = "empty"
            log.warning(f"Sin texto extraíble: {path.name}")
            return result
        chunks = ingest(
            source=f"onedrive/brand/{folder_origin}/{path.name}",
            text=text,
            metadata=metadata,
            namespace=BRAND_NAMESPACE,
            tenant=BRAND_TENANT,
        )
        result.update(status="ok", chunks=chunks)
        log.info(f"Ingestado: {path.name} → {chunks} chunks")

    elif ext in VISUAL_EXTENSIONS:
        # Assets visuales no tienen texto extraíble: indexamos metadatos descriptivos
        text = (
            f"Archivo visual de identidad de marca: {path.name}\n"
            f"Formato: {ext.lstrip('.')}\n"
            f"Carpeta de origen OneDrive: {folder_origin}\n"
            f"Ruta local: {path}\n"
            f"Organización: Reamerica\n"
            f"Categoría: identidad visual, brand assets\n"
        )
        chunks = ingest(
            source=f"onedrive/brand/{folder_origin}/{path.name}",
            text=text,
            metadata={**metadata, "is_visual": True},
            namespace=BRAND_NAMESPACE,
            tenant=BRAND_TENANT,
        )
        result.update(status="metadata_only", chunks=chunks)
        log.info(f"Metadata ingestada: {path.name}")

    return result


# ── Flujo principal ────────────────────────────────────────────────────────────

def run(user: str, site_id: str = None) -> list:
    from services.onedrive import onedrive_sync_folder, sharepoint_sync_folder

    all_results = []
    total_files = 0
    total_chunks = 0
    errors = []

    for folder_name in BRAND_FOLDERS:
        dest = DOWNLOAD_DIR / folder_name.replace(" ", "_")
        log.info(f"\n{'─'*60}")
        log.info(f"Carpeta: '{folder_name}' → {dest}")

        try:
            if site_id:
                files = sharepoint_sync_folder(
                    site_id=site_id, folder_path=folder_name,
                    dest_dir=str(dest), tenant=BRAND_TENANT,
                )
            else:
                files = onedrive_sync_folder(
                    folder_path=folder_name, dest_dir=str(dest),
                    user=user, tenant=BRAND_TENANT,
                )

            log.info(f"Archivos descargados: {len(files)}")
            total_files += len(files)

            for filepath in files:
                r = ingest_file(Path(filepath), folder_origin=folder_name)
                all_results.append(r)
                if r["status"] == "ok":
                    total_chunks += r["chunks"]
                elif r["status"] not in ("metadata_only", "empty", "skip"):
                    errors.append(r["file"])

        except Exception as e:
            log.error(f"Error procesando carpeta '{folder_name}': {e}")
            errors.append(f"CARPETA:{folder_name} — {e}")

    _print_report(all_results, total_files, total_chunks, errors)
    return all_results


def _print_report(results: list, total_files: int, total_chunks: int, errors: list):
    STATUS_ICON = {"ok": "✓", "metadata_only": "~", "empty": "○", "skip": "-"}
    print("\n" + "=" * 62)
    print("  REPORTE — Brand Assets Reamerica → RAG")
    print("=" * 62)
    print(f"  Fecha:               {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
    print(f"  Carpetas procesadas: {BRAND_FOLDERS}")
    print(f"  Archivos descargados:{total_files}")
    print(f"  Chunks ingestados:   {total_chunks}")
    print(f"  Errores:             {len(errors)}")
    print()
    print("  Detalle:")
    for r in results:
        icon = STATUS_ICON.get(r["status"], "?")
        print(f"    {icon} [{r['folder']}] {r['name']}"
              f"  →  {r['status']} ({r['chunks']} chunks)")
    if errors:
        print("\n  Errores:")
        for e in errors:
            print(f"    ! {e}")
    print("=" * 62)


def main():
    parser = argparse.ArgumentParser(description="Ingesta brand assets OneDrive → RAG")
    parser.add_argument("--user", default=os.environ.get("ONEDRIVE_REAMERICA_USER"),
                        help="Email del usuario OneDrive")
    parser.add_argument("--site", default=os.environ.get("ONEDRIVE_REAMERICA_SITE_ID"),
                        help="ID del site SharePoint (opcional)")
    parser.add_argument("--list-sites", action="store_true",
                        help="Lista los sites SharePoint disponibles y sale")
    args = parser.parse_args()

    # Carga el vault antes de cualquier operación de red
    try:
        from services.vault import load_all_to_env
        n = load_all_to_env()
        log.info(f"Vault: {n} secrets cargados al entorno")
    except Exception as e:
        log.warning(f"Vault no disponible (continuando con env vars): {e}")

    if args.list_sites:
        from services.onedrive import sharepoint_list_sites
        sites = sharepoint_list_sites(tenant=BRAND_TENANT)
        print(f"\n{len(sites)} site(s) SharePoint encontrados:")
        for s in sites:
            print(f"  id={s['id']!r:60s}  {s.get('displayName','?')}  {s.get('webUrl','')}")
        return

    if not args.user:
        print("ERROR: --user requerido (o ONEDRIVE_REAMERICA_USER en env/vault)")
        sys.exit(1)

    run(user=args.user, site_id=args.site)


if __name__ == "__main__":
    main()
