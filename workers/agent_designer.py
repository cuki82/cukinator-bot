"""
workers/agent_designer.py — Agent Designer (FastAPI service :3340).

Recibe pedidos de diseño del bot y genera assets corporativos:
  - HTML/CSS/Tailwind (landings, dashboards, componentes)
  - PDFs con branding (brochures, reportes, propuestas)
  - PPTX (presentaciones con layouts corporativos)
  - Critique (revisar diseño existente: usabilidad, a11y, consistencia)

RAG:
  <tenant>.kb_documents namespace 'brand' con el manual de identidad
  (logos, paleta, tipografías, plantillas). Si hay material ingestado,
  todo output lo aplica.

No genera imágenes raster propias — para eso hay que integrar DALL-E 3
o Stable Diffusion aparte (tool image_gen pendiente).

Arrancar con:
    uvicorn workers.agent_designer:app --host 0.0.0.0 --port 3340
"""
import os
import sys
import json
import logging
import tempfile
import time
from pathlib import Path
from typing import Optional, List

# Cargar vault
try:
    sys.path.insert(0, os.environ.get("REPO_PATH", "/home/cukibot/cukinator-bot"))
    from services.vault import load_all_to_env
    load_all_to_env()
except Exception:
    pass

import anthropic
from fastapi import FastAPI, HTTPException, Header
from pydantic import BaseModel

log = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")

app = FastAPI(title="Cukinator Agent Designer")

ANTHROPIC_KEY = os.environ.get("ANTHROPIC_KEY") or os.environ.get("ANTHROPIC_API_KEY", "")
DESIGNER_SECRET = os.environ.get("DESIGNER_SECRET") or os.environ.get("WORKER_SECRET", "cuki-designer-secret")
DEFAULT_MODEL = os.environ.get("DESIGNER_MODEL", "claude-sonnet-4-6")
OUTPUT_DIR = Path(os.environ.get("DESIGNER_OUTPUT", "/tmp/designer"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

client = anthropic.Anthropic(api_key=ANTHROPIC_KEY) if ANTHROPIC_KEY else None

log.info(
    "designer boot: ANTHROPIC_KEY=%s · DESIGNER_SECRET=%s · default_model=%s",
    "set" if ANTHROPIC_KEY else "MISSING",
    "set" if DESIGNER_SECRET else "MISSING",
    DEFAULT_MODEL,
)


# ── Models ────────────────────────────────────────────────────────────────

class DesignTask(BaseModel):
    task_id: str
    type: str              # 'html' | 'pdf' | 'pptx' | 'critique'
    tenant: str = "reamerica"
    chat_id: Optional[int] = None
    brief: str             # qué hay que diseñar
    target: Optional[str] = None      # audiencia / destino (ej. 'cliente', 'interno')
    reference: Optional[str] = None   # para critique: el HTML/texto a revisar
    model: Optional[str] = None       # override default


class DesignResult(BaseModel):
    task_id: str
    status: str            # 'ok' | 'error'
    type: str
    summary: str
    output_text: Optional[str] = None    # HTML / markdown / critique result
    output_file: Optional[str] = None    # path al PDF/PPTX generado
    brand_chunks_used: int = 0
    tokens_in: int = 0
    tokens_out: int = 0
    duration_s: float = 0
    errors: list = []


# ── Brand context via RAG ─────────────────────────────────────────────────

def _fetch_brand_context(tenant: str, brief: str, top_k: int = 5) -> str:
    """Busca chunks del manual de identidad del tenant. Namespace 'brand'."""
    try:
        from modules.rag_kb import build_context
        return build_context(
            query=f"identidad visual branding {brief}",
            top_k=top_k, namespace="brand",
            tenant=tenant, with_citations=False,
        ) or ""
    except Exception as e:
        log.debug(f"brand rag skip: {e}")
        return ""


# ── Generators por tipo ───────────────────────────────────────────────────

HTML_SYSTEM = """Sos un diseñador web senior. Generás HTML semántico + Tailwind CSS.

Reglas:
- HTML5 válido, responsive mobile-first, accesible WCAG 2.1 AA.
- Tailwind v3 utility-first. Sin CSS custom salvo cuando sea necesario.
- Jerarquía visual clara, contraste suficiente, espaciado consistente.
- Tipografía: sans-serif moderna (Inter / Geist / system-ui).
- Si te paso un manual de identidad, APLICALO estrictamente: colores exactos
  (hex), tipografías, proporciones, tono. El diseño debe verse del tenant.
- Output SOLO el HTML completo (<!DOCTYPE html>...<html>...). Sin comentarios
  ni explicaciones fuera del código. Una sola respuesta con el archivo listo."""

PDF_SYSTEM = """Sos un diseñador de presentaciones/PDFs corporativos. Tu output
es un plan estructurado que luego se renderea con fpdf2.

Formato de respuesta (JSON estricto, sin texto fuera):
{
  "title": "...",
  "subtitle": "...",
  "brand": {"primary_color": "#hex", "secondary_color": "#hex", "font": "Helvetica"},
  "pages": [
    {"heading": "...", "body": "markdown con **bold** y bullets", "layout": "cover|text|two_col|quote"}
  ]
}

Si te paso un manual de identidad, usá SUS colores y tipografías.
Tono directo, profesional, en español argentino."""

PPTX_SYSTEM = """Sos un diseñador de presentaciones corporativas. Tu output
es un plan JSON estructurado que se renderea con python-pptx.

Formato (JSON estricto, sin texto fuera ni triple-backticks):
{
  "title": "...",
  "subtitle": "...",
  "brand": {"primary_color": "#hex", "accent_color": "#hex", "font": "Calibri"},
  "slides": [
    {
      "layout": "title|section|content|two_col|quote|closing",
      "title": "...",
      "subtitle": "...",
      "bullets": ["...", "..."],
      "body_left": "...",
      "body_right": "...",
      "quote": "...",
      "attribution": "..."
    }
  ]
}

Reglas:
- Slide 1 SIEMPRE layout=title con title + subtitle.
- Slide final con layout=closing (call-to-action o cierre).
- 5-12 slides totales. Bullets ≤6 por slide, cada uno ≤14 palabras.
- Si hay manual de identidad, usá esos colores y fuente. Sino: primary='#1e3a5f' accent='#d97706' font='Calibri'.
- Tono directo, profesional, español argentino.
- Texto crudo, sin markdown ni emojis (los emojis rompen render con fuentes default)."""

CRITIQUE_SYSTEM = """Sos un crítico de diseño senior. Revisás HTML/mockups y
devolvés feedback accionable.

Estructura:
1. 🎯 Jerarquía visual — qué destaca, qué pierde el foco
2. 🎨 Consistencia de marca — alineación con el manual de identidad (si aplica)
3. ♿ Accesibilidad WCAG 2.1 AA — contraste, focus, aria, semántica
4. ✍️ UX writing — claridad, tono, microcopy
5. 📐 Layout y espaciado — grid, ritmo, respiración
6. 🔧 Acciones concretas — lista priorizada de cambios específicos

Directo, sin rodeos. No elogies por elogiar. Si algo está bien, decílo corto."""


def _run_claude(system: str, user: str, model: str) -> dict:
    """Llamada simple a Claude. Retorna dict con text, usage."""
    if not client:
        return {"error": "ANTHROPIC_KEY no configurado"}
    try:
        resp = client.messages.create(
            model=model,
            max_tokens=8000,
            system=system,
            messages=[{"role": "user", "content": user}],
        )
        text = ""
        for block in resp.content:
            if hasattr(block, "text"):
                text += block.text
        usage = getattr(resp, "usage", None)
        return {
            "text": text,
            "tokens_in":  getattr(usage, "input_tokens", 0) or 0 if usage else 0,
            "tokens_out": getattr(usage, "output_tokens", 0) or 0 if usage else 0,
        }
    except Exception as e:
        return {"error": str(e)}


def _gen_html(task: DesignTask, brand_ctx: str) -> dict:
    user = f"Brief: {task.brief}\n"
    if task.target:
        user += f"Audiencia/destino: {task.target}\n"
    if brand_ctx:
        user += f"\nManual de identidad del tenant:\n{brand_ctx}\n"
    user += "\nDevolvé el HTML completo."
    return _run_claude(HTML_SYSTEM, user, task.model or DEFAULT_MODEL)


def _gen_pdf_plan(task: DesignTask, brand_ctx: str) -> dict:
    user = f"Brief: {task.brief}\n"
    if task.target:
        user += f"Audiencia/destino: {task.target}\n"
    if brand_ctx:
        user += f"\nManual de identidad:\n{brand_ctx}\n"
    user += "\nDevolvé SOLO el JSON del plan del PDF."
    return _run_claude(PDF_SYSTEM, user, task.model or DEFAULT_MODEL)


def _render_pdf_from_plan(plan: dict, output_path: Path) -> None:
    """Usa fpdf2 para renderear el plan JSON al path dado."""
    from fpdf import FPDF
    brand = plan.get("brand", {})
    primary = brand.get("primary_color", "#1e3a5f").lstrip("#")
    try:
        pr, pg, pb = int(primary[:2], 16), int(primary[2:4], 16), int(primary[4:6], 16)
    except Exception:
        pr, pg, pb = 30, 58, 95

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)

    for page in plan.get("pages", []):
        layout = page.get("layout", "text")
        pdf.add_page()
        if layout == "cover":
            pdf.set_fill_color(pr, pg, pb)
            pdf.rect(0, 0, 210, 297, "F")
            pdf.set_text_color(255, 255, 255)
            pdf.set_font("Helvetica", "B", 24)
            pdf.set_y(100)
            pdf.multi_cell(0, 10, (page.get("heading") or plan.get("title") or "").encode("latin-1", errors="replace").decode("latin-1"))
            pdf.set_font("Helvetica", "", 14)
            pdf.multi_cell(0, 8, (page.get("body") or plan.get("subtitle") or "").encode("latin-1", errors="replace").decode("latin-1"))
        else:
            pdf.set_text_color(pr, pg, pb)
            pdf.set_font("Helvetica", "B", 16)
            heading = (page.get("heading") or "").encode("latin-1", errors="replace").decode("latin-1")
            pdf.cell(0, 10, heading, ln=1)
            pdf.ln(2)
            pdf.set_text_color(30, 30, 30)
            pdf.set_font("Helvetica", "", 11)
            body = (page.get("body") or "").encode("latin-1", errors="replace").decode("latin-1")
            for line in body.split("\n"):
                pdf.multi_cell(0, 6, line)
    pdf.output(str(output_path))


def _gen_pptx_plan(task: DesignTask, brand_ctx: str) -> dict:
    user = f"Brief: {task.brief}\n"
    if task.target:
        user += f"Audiencia/destino: {task.target}\n"
    if brand_ctx:
        user += f"\nManual de identidad:\n{brand_ctx}\n"
    user += "\nDevolvé SOLO el JSON del plan del PPTX."
    return _run_claude(PPTX_SYSTEM, user, task.model or DEFAULT_MODEL)


def _hex_to_rgb(hex_color: str, fallback=(30, 58, 95)):
    try:
        h = hex_color.lstrip("#")
        return int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return fallback


def _render_pptx_from_plan(plan: dict, output_path: Path) -> None:
    """python-pptx render. Crea slides según layouts del plan."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    brand = plan.get("brand", {})
    pr, pg, pb = _hex_to_rgb(brand.get("primary_color", "#1e3a5f"))
    ar, ag, ab = _hex_to_rgb(brand.get("accent_color", "#d97706"), fallback=(217, 119, 6))
    font_name = brand.get("font", "Calibri")
    primary = RGBColor(pr, pg, pb)
    accent = RGBColor(ar, ag, ab)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    dark = RGBColor(0x1E, 0x1E, 0x1E)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_text(slide, x, y, w, h, text, size=18, bold=False, color=dark, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = ""
        if not text:
            return tb
        for i, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
        return tb

    def fill_bg(slide, color):
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        bg.shadow.inherit = False
        # send back
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    def add_accent_bar(slide, color):
        bar = slide.shapes.add_shape(1, 0, Inches(7.2), prs.slide_width, Inches(0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

    for sl in plan.get("slides", []):
        layout = sl.get("layout", "content")
        slide = prs.slides.add_slide(blank)

        if layout == "title":
            fill_bg(slide, primary)
            add_text(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.5),
                     sl.get("title") or plan.get("title", ""), size=44, bold=True, color=white)
            add_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.0),
                     sl.get("subtitle") or plan.get("subtitle", ""), size=22, color=white)
            add_accent_bar(slide, accent)

        elif layout == "section":
            fill_bg(slide, primary)
            add_text(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.5),
                     sl.get("title", ""), size=40, bold=True, color=white)
            if sl.get("subtitle"):
                add_text(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.0),
                         sl.get("subtitle", ""), size=20, color=white)

        elif layout == "two_col":
            add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                     sl.get("title", ""), size=28, bold=True, color=primary)
            add_text(slide, Inches(0.6), Inches(1.6), Inches(6), Inches(5.5),
                     sl.get("body_left", ""), size=16)
            add_text(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(5.5),
                     sl.get("body_right", ""), size=16)
            add_accent_bar(slide, accent)

        elif layout == "quote":
            fill_bg(slide, RGBColor(0xF5, 0xF5, 0xF5))
            add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2.5),
                     '"' + (sl.get("quote", "") or "") + '"',
                     size=32, bold=True, color=primary, align=PP_ALIGN.CENTER)
            if sl.get("attribution"):
                add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.8),
                         "— " + sl["attribution"], size=18, color=dark, align=PP_ALIGN.CENTER)

        elif layout == "closing":
            fill_bg(slide, primary)
            add_text(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5),
                     sl.get("title", "Gracias"), size=44, bold=True, color=white)
            if sl.get("subtitle"):
                add_text(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(1.0),
                         sl["subtitle"], size=20, color=white)
            add_accent_bar(slide, accent)

        else:  # content (default)
            add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                     sl.get("title", ""), size=28, bold=True, color=primary)
            if sl.get("subtitle"):
                add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
                         sl["subtitle"], size=16, color=dark)
            bullets = sl.get("bullets") or []
            body_text = sl.get("body", "")
            y = Inches(1.9 if sl.get("subtitle") else 1.5)
            if bullets:
                bullet_text = "\n".join(f"• {b}" for b in bullets)
                add_text(slide, Inches(0.8), y, Inches(11.5), Inches(5.0),
                         bullet_text, size=18)
            elif body_text:
                add_text(slide, Inches(0.8), y, Inches(11.5), Inches(5.0),
                         body_text, size=16)
            add_accent_bar(slide, accent)

    prs.save(str(output_path))


def _gen_critique(task: DesignTask, brand_ctx: str) -> dict:
    if not task.reference:
        return {"error": "critique requiere 'reference' con el HTML/texto a revisar"}
    user = f"Diseño a revisar:\n```\n{task.reference[:12000]}\n```\n\nBrief original: {task.brief}\n"
    if brand_ctx:
        user += f"\nManual de identidad:\n{brand_ctx}\n"
    return _run_claude(CRITIQUE_SYSTEM, user, task.model or DEFAULT_MODEL)


# ── FastAPI endpoints ─────────────────────────────────────────────────────

@app.get("/health")
def health():
    return {
        "status": "ok",
        "anthropic": bool(ANTHROPIC_KEY),
        "default_model": DEFAULT_MODEL,
        "output_dir": str(OUTPUT_DIR),
        "types_supported": ["html", "pdf", "pptx", "critique"],
    }


@app.post("/design", response_model=DesignResult)
def design(task: DesignTask, x_designer_secret: str = Header(None)):
    if x_designer_secret != DESIGNER_SECRET:
        raise HTTPException(status_code=401, detail="Unauthorized")

    start = time.time()
    errors: list = []
    brand_ctx = _fetch_brand_context(task.tenant, task.brief)
    brand_chunks = brand_ctx.count("[C") if brand_ctx else 0

    out_text = None
    out_file = None
    tokens_in = 0
    tokens_out = 0

    try:
        if task.type == "html":
            r = _gen_html(task, brand_ctx)
            if "error" in r:
                errors.append(r["error"])
            else:
                out_text = r["text"]
                tokens_in = r["tokens_in"]
                tokens_out = r["tokens_out"]
                # Guardar también como archivo
                out_file = str(OUTPUT_DIR / f"{task.task_id}.html")
                Path(out_file).write_text(out_text, encoding="utf-8")

        elif task.type == "pdf":
            r = _gen_pdf_plan(task, brand_ctx)
            if "error" in r:
                errors.append(r["error"])
            else:
                tokens_in = r["tokens_in"]
                tokens_out = r["tokens_out"]
                try:
                    plan = json.loads(r["text"].strip().strip("`").replace("json\n", "", 1))
                    out_file = str(OUTPUT_DIR / f"{task.task_id}.pdf")
                    _render_pdf_from_plan(plan, Path(out_file))
                    out_text = json.dumps(plan, indent=2, ensure_ascii=False)
                except Exception as ex:
                    errors.append(f"pdf render error: {ex}")
                    out_text = r["text"]  # devolver texto raw como fallback

        elif task.type == "critique":
            r = _gen_critique(task, brand_ctx)
            if "error" in r:
                errors.append(r["error"])
            else:
                out_text = r["text"]
                tokens_in = r["tokens_in"]
                tokens_out = r["tokens_out"]

        elif task.type == "pptx":
            r = _gen_pptx_plan(task, brand_ctx)
            if "error" in r:
                errors.append(r["error"])
            else:
                tokens_in = r["tokens_in"]
                tokens_out = r["tokens_out"]
                try:
                    raw = r["text"].strip()
                    # quitar fences ```json ... ```
                    if raw.startswith("```"):
                        raw = raw.split("```", 2)[1]
                        if raw.startswith("json"):
                            raw = raw[4:]
                        raw = raw.rstrip("`").strip()
                    plan = json.loads(raw)
                    out_file = str(OUTPUT_DIR / f"{task.task_id}.pptx")
                    _render_pptx_from_plan(plan, Path(out_file))
                    out_text = json.dumps(plan, indent=2, ensure_ascii=False)
                except Exception as ex:
                    errors.append(f"pptx render error: {ex}")
                    out_text = r["text"]

        else:
            errors.append(f"type desconocido: {task.type}")

    except Exception as e:
        errors.append(str(e))

    summary = f"designer {task.type} · {len(out_text or '')} chars"
    if out_file:
        summary += f" · file={Path(out_file).name}"

    return DesignResult(
        task_id=task.task_id,
        status="ok" if not errors else ("partial" if (out_text or out_file) else "error"),
        type=task.type,
        summary=summary,
        output_text=out_text,
        output_file=out_file,
        brand_chunks_used=brand_chunks,
        tokens_in=tokens_in,
        tokens_out=tokens_out,
        duration_s=round(time.time() - start, 1),
        errors=errors,
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=3340, log_level="info")
